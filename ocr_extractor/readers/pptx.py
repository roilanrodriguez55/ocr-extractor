"""PPTX reader — extract text directly from PowerPoint 2007+ files.

Each slide is wrapped in ``=== PAGE N ===`` / ``=== END PAGE N ===``
markers, matching the format produced by the PDF and multi-page TIFF
readers. Text is read from every shape that has a ``text_frame`` and
from tables (each row becomes one line).
"""

import re

from pptx import Presentation
from ocr_extractor.result import wrap_page


def _light_clean(text):
    """Collapse internal whitespace and strip the edges."""
    return re.sub(r"\s+", " ", text).strip()


def _shape_text_lines(shape):
    """Yield non-empty text lines from a shape.

    Handles:
    - Shapes with a ``text_frame`` (text boxes, placeholders, etc.).
    - Tables (each row becomes ``" | col1 | col2 | col3"``).
    """
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            text = _light_clean(para.text)
            if text:
                yield text
        return

    if shape.has_table:
        for row in shape.table.rows:
            cells = [_light_clean(cell.text) for cell in row.cells]
            line = " | " + " | ".join(cells)
            stripped = line.strip(" |")
            if stripped:
                yield stripped


def read_pptx(path, *, dpi=300, lang="eng", verbose=True):
    """Extract text from a ``.pptx`` file.

    Each slide produces one marker block. The text inside is built by
    walking every shape on the slide (text frames and tables) and
    emitting one line per paragraph or row.

    Parameters
    ----------
    path : str or pathlib.Path
        Path to the ``.pptx`` file.
    dpi : int, optional
        Unused; kept for signature uniformity. Defaults to ``300``.
    lang : str, optional
        Unused; kept for signature uniformity. Defaults to ``"eng"``.
    verbose : bool, optional
        If ``True``, print a per-slide progress line. Defaults to
        ``True``.

    Returns
    -------
    str
        Extracted text with ``=== PAGE N ===`` / ``=== END PAGE N ===``
        markers, one block per slide.
    """
    if verbose:
        print(f"Reading {path}...")

    prs = Presentation(path)
    all_text = ""

    for i, slide in enumerate(prs.slides, 1):
        if verbose:
            print(f"Processing slide {i}...")

        lines = []
        for shape in slide.shapes:
            lines.extend(_shape_text_lines(shape))

        body = "\n".join(lines)
        all_text += wrap_page(i, body)

    return all_text


__all__ = ["read_pptx"]