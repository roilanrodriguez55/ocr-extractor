"""Tests for the PPTX reader."""

import pytest
from pptx import Presentation

from ocr_extractor.readers.pptx import read_pptx


# ---------- synthetic PPTX fixtures ----------------------------------------


def _make_pptx(path, slides=None):
    """Build a PPTX where each slide has a list of text lines.

    Parameters
    ----------
    path : pathlib.Path
        Where to save the .pptx.
    slides : list[list[str]] or None
        One inner list per slide; each string becomes a paragraph on a
        text box added to that slide.
    """
    prs = Presentation()
    for slide_lines in slides or []:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
        left, top, width, height = 100, 100, 400, 300
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        for i, line in enumerate(slide_lines):
            if i == 0:
                tf.text = line
            else:
                p = tf.add_paragraph()
                p.text = line
    prs.save(str(path))
    return path


# ---------- tests ----------------------------------------------------------


class TestReadPptx:
    def test_empty_presentation_returns_markers_only(self, tmp_workdir):
        path = _make_pptx(tmp_workdir / "empty.pptx", slides=[])
        text = read_pptx(path, verbose=False)
        # No slides → no markers at all.
        assert text == ""

    def test_single_slide_has_one_marker_block(self, tmp_workdir):
        path = _make_pptx(
            tmp_workdir / "single.pptx",
            slides=[["Hello world"]],
        )
        text = read_pptx(path, verbose=False)
        assert "=== PAGE 1 ===" in text
        assert "=== END PAGE 1 ===" in text
        assert "Hello world" in text
        assert "=== PAGE 2 ===" not in text

    def test_multiple_slides_each_marked(self, tmp_workdir):
        path = _make_pptx(
            tmp_workdir / "multi.pptx",
            slides=[
                ["First slide content"],
                ["Second slide content"],
                ["Third slide content"],
            ],
        )
        text = read_pptx(path, verbose=False)
        for i in (1, 2, 3):
            assert f"=== PAGE {i} ===" in text
            assert f"=== END PAGE {i} ===" in text

        # Marker order matches slide order.
        assert text.index("PAGE 1") < text.index("PAGE 2") < text.index("PAGE 3")
        # Content order matches slide order.
        assert text.index("First slide") < text.index("Second slide") < text.index("Third slide")

    def test_preserves_punctuation(self, tmp_workdir):
        path = _make_pptx(
            tmp_workdir / "punct.pptx",
            slides=[["Hello, world! How are you?"]],
        )
        text = read_pptx(path, verbose=False)
        assert "," in text
        assert "!" in text
        assert "?" in text

    def test_drops_empty_paragraphs(self, tmp_workdir):
        # Build manually so we can include blank paragraphs.
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        textbox = slide.shapes.add_textbox(100, 100, 400, 300)
        tf = textbox.text_frame
        tf.text = "Real content"
        tf.add_paragraph().text = ""
        tf.add_paragraph().text = "   "
        tf.add_paragraph().text = "More content"
        prs.save(str(tmp_workdir / "with_empty.pptx"))

        text = read_pptx(tmp_workdir / "with_empty.pptx", verbose=False)
        assert "Real content" in text
        assert "More content" in text
        # Only the two non-empty paragraphs survive (filtering out the
        # marker lines).
        content_lines = [
            line
            for line in text.split("\n")
            if line.strip() and not line.startswith("===")
        ]
        assert content_lines == ["Real content", "More content"]

    def test_extracts_table_rows(self, tmp_workdir):
        """Tables on slides should yield one line per row."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        rows, cols = 2, 3
        table_shape = slide.shapes.add_table(rows, cols, 100, 100, 400, 100)
        for i in range(rows):
            for j in range(cols):
                table_shape.table.cell(i, j).text = f"R{i}C{j}"
        prs.save(str(tmp_workdir / "table.pptx"))

        text = read_pptx(tmp_workdir / "table.pptx", verbose=False)
        assert "R0C0 | R0C1 | R0C2" in text
        assert "R1C0 | R1C1 | R1C2" in text

    def test_combines_text_frames_and_tables(self, tmp_workdir):
        """Shapes of different types on the same slide are all visited."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # First: a text box.
        textbox = slide.shapes.add_textbox(100, 100, 400, 50)
        textbox.text_frame.text = "Title text"

        # Second: a table.
        table_shape = slide.shapes.add_table(1, 2, 100, 200, 400, 50)
        table_shape.table.cell(0, 0).text = "A"
        table_shape.table.cell(0, 1).text = "B"

        prs.save(str(tmp_workdir / "combo.pptx"))

        text = read_pptx(tmp_workdir / "combo.pptx", verbose=False)
        assert "Title text" in text
        assert "A | B" in text

    def test_verbose_prints_per_slide(self, tmp_workdir, capsys):
        path = _make_pptx(
            tmp_workdir / "v.pptx",
            slides=[["x"], ["y"]],
        )
        read_pptx(path, verbose=True)
        captured = capsys.readouterr()
        assert "Processing slide 1" in captured.out
        assert "Processing slide 2" in captured.out

    def test_missing_file_raises(self, tmp_workdir):
        with pytest.raises(FileNotFoundError):
            read_pptx(tmp_workdir / "nope.pptx", verbose=False)